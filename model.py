import os
import asyncio
import time
import sys
import subprocess
import shutil
from pathlib import Path
from PyPDF2 import PdfReader
from pptx import Presentation
from edge_tts import Communicate
from typing import List, Callable

if sys.platform == "win32":
    import pydub.utils
    import ctypes.wintypes
    import ctypes

    def _run_command_hide_console(args, **kwargs):
        CREATE_NO_WINDOW = 0x08000000
        kwargs["creationflags"] = CREATE_NO_WINDOW
        return subprocess.Popen(args, **kwargs)

    pydub.utils.Popen = _run_command_hide_console

    def obtener_escritorio_usuario():
        CSIDL_DESKTOP = 0
        SHGFP_TYPE_CURRENT = 0
        buf = ctypes.create_unicode_buffer(ctypes.wintypes.MAX_PATH)
        ctypes.windll.shell32.SHGetFolderPathW(None, CSIDL_DESKTOP, None, SHGFP_TYPE_CURRENT, buf)
        return buf.value
    
    def ocultar_archivo_windows(ruta):
        FILE_ATTRIBUTE_HIDDEN = 0x02
        ctypes.windll.kernel32.SetFileAttributesW(ruta, FILE_ATTRIBUTE_HIDDEN)

else:
    def obtener_escritorio_usuario():
        return str(Path.home() / "Desktop")
    
    def ocultar_archivo_windows(ruta):
        pass

def texto_a_fragmentos(texto: str, max_palabras: int = 350) -> List[str]:
    palabras = texto.split()
    return [" ".join(palabras[i:i+max_palabras]) for i in range(0, len(palabras), max_palabras)]

async def generar_fragmento_tts(texto: str, archivo_salida: str, voz: str):
    comunicador = Communicate(texto, voice=voz)
    await comunicador.save(archivo_salida)

def texto_a_voz_edge_fragmentos(
    fragmentos: List[str], 
    carpeta_salida: str, 
    voz: str, 
    log_func: Callable[[str], None], 
    actualizar_progreso_func: Callable[[int, int], None]
) -> List[str]:
    if not os.path.exists(carpeta_salida):
        os.makedirs(carpeta_salida)
        ocultar_archivo_windows(carpeta_salida)
        
    archivos_mp3 = []
    total = len(fragmentos)
    log_func(f"\n[INFO] Total fragmentos a procesar: {total}\n")
    
    for i, fragmento in enumerate(fragmentos):
        if not fragmento.strip():
            continue
            
        archivo_mp3 = os.path.join(carpeta_salida, f"parte_{i+1}.mp3")
        log_func(f"Procesando fragmento {i+1} de {total}...")
        
        try:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            loop.run_until_complete(generar_fragmento_tts(fragmento, archivo_mp3, voz=voz))
            
            archivos_mp3.append(archivo_mp3)
            actualizar_progreso_func(i + 1, total)
            time.sleep(0.5)
            
        except Exception as e:
            log_func(f"[ERROR] Fragmento {i+1} falló: {e}")
            break
            
    log_func(f"\n[RESUMEN] Se generaron {len(archivos_mp3)} archivos de audio.")
    return archivos_mp3

def leer_pdf_por_paginas(ruta_pdf: str) -> List[str]:
    try:
        reader = PdfReader(ruta_pdf)
        return [page.extract_text() or "" for page in reader.pages]
    except Exception:
        return []

def leer_pptx_por_diapositivas(ruta_pptx: str) -> List[str]:
    textos_diapositivas = []
    try:
        prs = Presentation(ruta_pptx)
        for slide in prs.slides:
            texto = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
            textos_diapositivas.append(texto.strip())
    except Exception:
        textos_diapositivas = []
    return textos_diapositivas

def leer_txt(ruta_txt: str) -> str:
    try:
        with open(ruta_txt, "r", encoding="utf-8") as f:
            return f.read()
    except Exception:
        return ""

def verificar_ffmpeg() -> bool:
    try:
        result = subprocess.run(['ffmpeg', '-version'], capture_output=True, check=False)
        return result.returncode == 0
    except Exception:
        return False

def instalar_ffmpeg():
    if sys.platform == "win32":
        try:
            subprocess.Popen(['winget', 'install', '--id', 'Gyan.FFmpeg', '-e', '--source', 'winget'])
            return True
        except Exception:
            return False
    return False

def eliminar_carpeta_temporal(carpeta: str):
    try:
        if os.path.exists(carpeta):
            shutil.rmtree(carpeta)
    except Exception as e:
        print(f"No se pudo eliminar la carpeta temporal: {e}")

class TTSModel:
    def __init__(self):
        self.proceso_activo = False
        self.voz_seleccionada = "es-CL-CatalinaNeural"
        self.archivos_generados = []
        self.ruta_archivo_actual = None
        self.carpeta_temporal = None
        
    def set_voz(self, voz: str):
        self.voz_seleccionada = voz

    def set_ruta_archivo(self, ruta: str):
        self.ruta_archivo_actual = ruta
        
    def iniciar_proceso(self, ruta_archivo: str):
        self.proceso_activo = True
        self.ruta_archivo_actual = ruta_archivo
        self.carpeta_temporal = os.path.join(obtener_escritorio_usuario(), "audios_temp_hidden")
        eliminar_carpeta_temporal(self.carpeta_temporal)

    def cancelar_proceso(self):
        self.proceso_activo = False

    def finalizar_proceso(self):
        self.proceso_activo = False
        eliminar_carpeta_temporal(self.carpeta_temporal)

    def mover_archivos_a_destino(self, carpeta_destino: str):
        for archivo in self.archivos_generados:
            nombre = os.path.basename(archivo)
            destino = os.path.join(carpeta_destino, nombre)
            shutil.move(archivo, destino)